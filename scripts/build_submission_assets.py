"""Build final defense assets from real outputs on the student's machine.

This script reads outputs/run_summary.json and outputs/db/violations.db, then
creates:
- outputs/final_assets/final_metrics.json
- outputs/final_assets/final_metrics.md
- outputs/final_assets/screenshots/*.jpg
- outputs/final_assets/demo_backup.mp4
- outputs/final_assets/final_report.docx
- outputs/final_assets/final_slide_deck.pptx

Run after the terminal or Streamlit demo has produced real outputs.
"""

from __future__ import annotations

import argparse
import json
import shutil
import sqlite3
from datetime import datetime
from pathlib import Path

import cv2


def load_json(path: Path) -> dict:
    return json.loads(path.read_text(encoding="utf-8")) if path.exists() else {}


def load_violations(db_path: Path) -> list[dict]:
    if not db_path.exists():
        return []
    conn = sqlite3.connect(db_path)
    conn.row_factory = sqlite3.Row
    try:
        return [dict(r) for r in conn.execute("SELECT * FROM violations ORDER BY id ASC").fetchall()]
    finally:
        conn.close()


def extract_screenshots(video_path: Path, out_dir: Path, max_images: int = 6) -> list[str]:
    out_dir.mkdir(parents=True, exist_ok=True)
    if not video_path.exists():
        return []
    cap = cv2.VideoCapture(str(video_path))
    if not cap.isOpened():
        return []
    total = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
    if total <= 0:
        frame_ids = [0]
    else:
        frame_ids = sorted(set(int(total * p) for p in [0.10, 0.25, 0.40, 0.55, 0.70, 0.85]))[:max_images]
    paths = []
    for idx, frame_id in enumerate(frame_ids, 1):
        cap.set(cv2.CAP_PROP_POS_FRAMES, frame_id)
        ok, frame = cap.read()
        if ok:
            path = out_dir / f"screenshot_{idx:02d}.jpg"
            cv2.imwrite(str(path), frame)
            paths.append(str(path))
    cap.release()
    return paths


def write_metrics_md(metrics: dict, path: Path) -> None:
    lines = [
        "# Final real metrics", "",
        f"Generated at: {metrics['generated_at']}", "",
        "| Item | Value |",
        "|---|---:|",
        f"| Processed frames | {metrics.get('processed_frames', 0)} |",
        f"| Average processing FPS | {metrics.get('average_processing_fps', 0)} |",
        f"| Confirmed violation events | {metrics.get('violation_count', 0)} |",
        f"| Confirmed no-helmet rider tracks | {metrics.get('confirmed_no_helmet_tracks', 0)} |",
        f"| OCR readable plates | {metrics.get('ocr_readable_count', 0)} |",
        f"| OCR unreadable plates | {metrics.get('ocr_unreadable_count', 0)} |",
        "", "## Detection counts", "",
    ]
    for name, count in metrics.get("detection_counts", {}).items():
        lines.append(f"- {name}: {count}")
    lines.extend(["", "## Violation records", ""])
    for v in metrics.get("violations", []):
        lines.append(
            f"- ID {v.get('id')}: Track {v.get('track_id')}, conf={v.get('confidence'):.3f}, "
            f"plate={v.get('plate_text')}, image={v.get('image_path')}"
        )
    path.write_text("\n".join(lines), encoding="utf-8")


def build_report_docx(metrics: dict, out_path: Path) -> None:
    from docx import Document
    from docx.shared import Inches

    doc = Document()
    doc.add_heading("Motorcycle Helmet Violation Detection", 0)
    doc.add_paragraph("Final MVP report generated from real local demo outputs.")
    doc.add_heading("1. Scope", 1)
    doc.add_paragraph(
        "The system detects motorcycle helmet violations from MP4/webcam video using YOLOv8s, "
        "ByteTrack, temporal confirmation, SQLite logging and optional OCR."
    )
    doc.add_heading("2. Pipeline", 1)
    doc.add_paragraph(
        "Input video -> OpenCV frames -> YOLO detection -> ByteTrack ID -> no_helmet temporal confirmation -> "
        "full-frame evidence -> SQLite/CSV -> annotated video/Streamlit; optional OCR uses a heuristic crop."
    )
    doc.add_heading("3. Real demo metrics", 1)
    table = doc.add_table(rows=1, cols=2)
    table.style = "Table Grid"
    table.rows[0].cells[0].text = "Metric"
    table.rows[0].cells[1].text = "Value"
    rows = [
        ("Processed frames", metrics.get("processed_frames", 0)),
        ("Average processing FPS", metrics.get("average_processing_fps", 0)),
        ("Confirmed violations", metrics.get("violation_count", 0)),
        ("Confirmed no-helmet rider tracks", metrics.get("confirmed_no_helmet_tracks", 0)),
        ("OCR readable plates", metrics.get("ocr_readable_count", 0)),
    ]
    for k, v in rows:
        cells = table.add_row().cells
        cells[0].text = str(k)
        cells[1].text = str(v)

    doc.add_heading("4. Evidence screenshots", 1)
    for image in metrics.get("evidence_images", [])[:4]:
        p = Path(image)
        if p.exists():
            doc.add_paragraph(p.name)
            doc.add_picture(str(p), width=Inches(5.6))

    doc.add_heading("5. Limitations", 1)
    doc.add_paragraph(
        "OCR can return UNREAD when the license plate is too small, blurred or not associated with the rider. "
        "The system is a classroom decision-support MVP, not an autonomous legal enforcement system."
    )
    doc.save(out_path)


def build_pptx(metrics: dict, out_path: Path) -> None:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    def title_slide(title: str, subtitle: str = ""):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(0.7), Inches(0.8), Inches(12), Inches(1.2))
        p = box.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(40)
        p.font.bold = True
        if subtitle:
            s = slide.shapes.add_textbox(Inches(0.75), Inches(2.0), Inches(11.5), Inches(0.6))
            sp = s.text_frame.paragraphs[0]
            sp.text = subtitle
            sp.font.size = Pt(22)
        return slide

    def bullets(title: str, items: list[str]):
        slide = title_slide(title)
        box = slide.shapes.add_textbox(Inches(0.9), Inches(1.8), Inches(11.7), Inches(4.8))
        tf = box.text_frame
        tf.clear()
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = item
            p.font.size = Pt(23)
            p.space_after = Pt(10)
        return slide

    title_slide("Helmet Violation Detection MVP", "YOLOv8s • ByteTrack • SQLite • Streamlit")
    bullets("Problem & Scope", [
        "Detect riders without helmets from video.",
        "Prioritize a working ML/CV pipeline over incomplete infrastructure.",
        "RTSP/React/Django/Stringee are future work, not core rubric requirements.",
    ])
    bullets("Pipeline", [
        "MP4/Webcam -> OpenCV frames -> YOLOv8s detection",
        "Confidence filtering + NMS -> ByteTrack Track ID",
        "no_helmet confirmed across multiple frames -> full-frame evidence",
        "SQLite/CSV evidence -> annotated video/Streamlit; OCR is optional and experimental",
    ])
    bullets("Why YOLOv8s", [
        "One-stage detector: bounding box and class in one forward pass.",
        "The small model balances detector accuracy and demo speed, and can still run on CPU.",
        "Faster R-CNN is a valid baseline but heavier for a realtime demo.",
    ])
    bullets("Tracking & Violation Rule", [
        "ByteTrack assigns persistent IDs across frames.",
        "A violation is saved only after the same no_helmet Track ID appears repeatedly.",
        "This reduces single-frame false positives and duplicate records.",
    ])
    bullets("Real Demo Metrics", [
        f"Processed frames: {metrics.get('processed_frames', 0)}",
        f"Average processing FPS: {metrics.get('average_processing_fps', 0)}",
        f"Confirmed violations: {metrics.get('violation_count', 0)}",
        f"Confirmed no-helmet rider tracks: {metrics.get('confirmed_no_helmet_tracks', 0)}",
        f"OCR readable plates: {metrics.get('ocr_readable_count', 0)}",
    ])
    bullets("Evidence Storage", [
        "Full-frame evidence image is saved for human review.",
        "Plate crop is saved separately for OCR debugging.",
        "SQLite stores time, Track ID, confidence, plate text and image paths.",
    ])
    bullets("Limitations", [
        "Small/blurred plates can produce UNREAD OCR.",
        "The bundled model has no plate class; optional OCR uses a heuristic lower-rider crop.",
        "The output requires human review and is not legal evidence by itself.",
    ])
    bullets("Future Work", [
        "Train on larger Vietnam-specific traffic data.",
        "Use a relationship model for rider-motorcycle-plate grouping.",
        "Add RTSP, reviewer dashboard, authentication and deployment after the AI core is stable.",
    ])
    title_slide("Demo", "Open the annotated MP4, SQLite/CSV records and Streamlit UI.")
    prs.save(out_path)


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("--outputs", default="outputs")
    parser.add_argument("--video", default="")
    args = parser.parse_args()

    outputs = Path(args.outputs)
    final_dir = outputs / "final_assets"
    final_dir.mkdir(parents=True, exist_ok=True)
    screen_dir = final_dir / "screenshots"

    summary = load_json(outputs / "run_summary.json")
    db_path = Path(summary.get("db_path", outputs / "db" / "violations.db"))
    violations = load_violations(db_path)

    video_path = Path(args.video) if args.video else Path(summary.get("output_video", outputs / "videos" / "gpu_result.mp4"))
    if video_path.exists():
        shutil.copy2(video_path, final_dir / "demo_backup.mp4")

    screenshots = extract_screenshots(video_path, screen_dir)
    evidence_images = [v.get("image_path", "") for v in violations if v.get("image_path")]
    plate_images = [v.get("plate_image_path", "") for v in violations if v.get("plate_image_path")]

    readable = [v for v in violations if v.get("plate_text") and v.get("plate_text") != "UNREAD"]
    unreadable = [v for v in violations if not v.get("plate_text") or v.get("plate_text") == "UNREAD"]

    metrics = {
        "generated_at": datetime.now().isoformat(timespec="seconds"),
        **summary,
        "violation_count": len(violations),
        "violations": violations,
        "ocr_readable_count": len(readable),
        "ocr_unreadable_count": len(unreadable),
        "screenshots": screenshots,
        "evidence_images": evidence_images,
        "plate_images": plate_images,
        "demo_backup": str(final_dir / "demo_backup.mp4") if video_path.exists() else "",
    }

    (final_dir / "final_metrics.json").write_text(json.dumps(metrics, indent=2, ensure_ascii=False), encoding="utf-8")
    write_metrics_md(metrics, final_dir / "final_metrics.md")
    build_report_docx(metrics, final_dir / "final_report.docx")
    build_pptx(metrics, final_dir / "final_slide_deck.pptx")

    print("Final assets created:")
    for p in [
        final_dir / "final_metrics.json",
        final_dir / "final_metrics.md",
        final_dir / "demo_backup.mp4",
        final_dir / "final_report.docx",
        final_dir / "final_slide_deck.pptx",
    ]:
        print(" -", p)
    print("Screenshots:", screen_dir)


if __name__ == "__main__":
    main()
